#!/usr/bin/env python3
"""
ingest.py

CLI & library to ingest course materials (PDF, PPTX, MD/TXT)
-> chunk -> embed -> insert into Chroma.

Usage (CLI):
  python ingest.py --input path/to/file_or_dir \
                   --collection course_calculus \
                   --course-id calculus101 \
                   --lecture-id lecture_02 \
                   --persist-dir ./chroma_db

You can also import functions:
  from ingest import ingest_files
  ingest_files([...], collection_name="course_x", ...)
"""

import os
import argparse
import uuid
import logging
from pathlib import Path
from typing import List, Dict, Tuple, Optional

import chromadb
from chromadb.config import Settings
from tqdm import tqdm

# Text extraction
from PyPDF2 import PdfReader
from pptx import Presentation

# Embeddings & tokenization
import tiktoken
from groq.types import CreateEmbeddingResponse, Embedding
import os


# Optional OpenAI / fallback transformer
try:
    import openai
except Exception:
    openai = None

try:
    from sentence_transformers import SentenceTransformer
except Exception:
    SentenceTransformer = None

# ---------- Configuration ----------
DEFAULT_EMBEDDING_MODEL = "text-embedding-3-small"  # OpenAI
# token encoding choice for counting (cl100k_base works for OpenAI's models)
ENCODING_NAME = "cl100k_base"

# ---------- Logging ----------
logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
logger = logging.getLogger("ingest")

# ---------- Utilities: text extraction ----------
def extract_text_from_pdf(path: str) -> List[Tuple[str, Optional[int]]]:
    """Return list of (text, page_number) for each page."""
    reader = PdfReader(path)
    pages = []
    print(f"PDF has {len(reader.pages)} pages")
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        pages.append((text, i + 1))
    return pages


def extract_text_from_pptx(path: str) -> List[Tuple[str, Optional[int]]]:
    """Return list of (text, slide_index) for each slide."""
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
        slides.append(("\n".join(parts), i + 1))
    return slides


def extract_text_from_txt_or_md(path: str) -> List[Tuple[str, Optional[int]]]:
    """Read entire file as a single chunk (page=None)."""
    with open(path, "r", encoding="utf-8") as f:
        return [(f.read(), None)]


# ---------- Token-aware chunking ----------
def get_token_encoder():
    return tiktoken.get_encoding(ENCODING_NAME)


def chunk_text_tokenwise(text: str, max_tokens: int = 500, overlap: int = 100) -> List[str]:
    """
    Token-slice the text into chunks using tiktoken encoding and decode back.
    This ensures chunks are close to token budget.
    """
    enc = get_token_encoder()
    tokens = enc.encode(text)
    if not tokens:
        return []
    chunks = []
    step = max_tokens - overlap
    i = 0
    while i < len(tokens):
        chunk_tokens = tokens[i : i + max_tokens]
        chunk_text = enc.decode(chunk_tokens)
        chunks.append(chunk_text)
        i += step
    return chunks


# ---------- Embeddings ----------
class Embedder:
    def __init__(self, use_openai: bool = True, openai_api_key: Optional[str] = None):
        self.use_openai = use_openai and (openai is not None) and bool(openai_api_key)
        if self.use_openai:
            openai.api_key = openai_api_key
            logger.info("Using OpenAI embeddings")
        else:
            if SentenceTransformer is None:
                raise RuntimeError(
                    "OpenAI key not found and sentence-transformers not installed. "
                    "Install sentence-transformers or set OPENAI_API_KEY."
                )
            logger.info("Using sentence-transformers fallback for embeddings")
            self.model = SentenceTransformer("all-MiniLM-L6-v2")

    def embed_many(self, texts: List[str]) -> List[List[float]]:
        if self.use_openai:
            # OpenAI accepts batch up to certain size; do small batches
            results = []
            BATCH = 16
            for i in range(0, len(texts), BATCH):
                batch = texts[i : i + BATCH]
                resp = openai.Embedding.create(model=DEFAULT_EMBEDDING_MODEL, input=batch)
                # resp["data"] is a list with embedding per input
                for d in resp["data"]:
                    results.append(d["embedding"])
            return results
        else:
            embeddings = self.model.encode(texts, show_progress_bar=False, convert_to_numpy=False)
            # If tensorflow/torch returns numpy arrays, convert to list
            return [e.tolist() for e in embeddings]


# ---------- Chroma indexing ----------
def get_chroma_client(persist_directory: Optional[str] = None) -> chromadb.Client:
    """
    Create a Chroma client. For local persisted DB we use duckdb+parquet via Settings.
    If persist_directory not given, uses ephemeral in-memory Chroma.
    """
    if persist_directory:
        settings = Settings(chroma_db_impl="duckdb+parquet", persist_directory=persist_directory)
        client = chromadb.Client(settings)
    else:
        client = chromadb.Client()
    return client


def ensure_collection(client: chromadb.Client, collection_name: str):
    """Get or create collection."""
    try:
        # coll = client.get_collection(collection_name)
        print("creating collection")
        
        return client.create_collection(name=collection_name)
        # return coll
    except Exception:
        return client.create_collection(name=collection_name)


def insert_chunks_to_collection(
    coll,
    ids: List[str],
    documents: List[str],
    metadatas: List[Dict],
    embeddings: Optional[List[List[float]]] = None,
):
    """
    Insert vectors and metadata into Chroma collection.
    """
    if embeddings is not None:
        coll.add(ids=ids, documents=documents, metadatas=metadatas, embeddings=embeddings)
    else:
        coll.add(ids=ids, documents=documents, metadatas=metadatas)


# ---------- Main ingestion flow ----------
def process_single_file(
    path: str,
    course_id: str,
    lecture_id: Optional[str],
    embedder: Embedder,
    coll,
    max_tokens: int = 500,
    overlap: int = 100,
):
    """
    Extracts text from one file, chunks it, builds embeddings, then inserts into collection.
    Returns number of chunks indexed.
    """
    path = Path(path)
    suffix = path.suffix.lower()
    logger.info(f"Processing {path} (type={suffix})")

    if suffix == ".pdf":
        parts = extract_text_from_pdf(str(path))  # list[(text, page)]
    elif suffix in [".pptx"]:
        parts = extract_text_from_pptx(str(path))  # list[(text, slide)]
    elif suffix in [".md", ".txt"]:
        parts = extract_text_from_txt_or_md(str(path))
    else:
        logger.warning(f"Skipping unsupported file type: {path}")
        return 0

    all_chunks = []
    all_metadatas = []
    all_ids = []

    for (text, page_idx) in parts:
        if not text or len(text.strip()) == 0:
            continue
        chunks = chunk_text_tokenwise(text, max_tokens=max_tokens, overlap=overlap)
        for i, chunk in enumerate(chunks):
            uid = str(uuid.uuid4())
            metadata = {
                "course_id": course_id,
                "lecture_id": lecture_id,
                "source": path.name,
                "page_or_slide": page_idx,
                "chunk_index": i,
            }
            all_chunks.append(chunk)
            all_metadatas.append(metadata)
            all_ids.append(uid)

    if not all_chunks:
        logger.info("No chunks extracted, skipping embedding/insert.")
        return 0

    logger.info(f"Generating embeddings for {len(all_chunks)} chunks...")
    embeddings = embedder.embed_many(all_chunks)

    logger.info(f"Inserting {len(all_chunks)} chunks to Chroma collection...")
    insert_chunks_to_collection(coll, all_ids, all_chunks, all_metadatas, embeddings)
    return len(all_chunks)


def ingest_files(
    paths: List[str],
    collection_name: str,
    course_id: str,
    lecture_id: Optional[str],
    persist_directory: Optional[str],
    openai_api_key: Optional[str],
    max_tokens: int = 500,
    overlap: int = 100,
):
    """
    Top-level ingestion: connect to chroma, create collection, and process each file.
    """
    client = get_chroma_client(persist_directory=persist_directory)
    try:
        coll = client.get_or_create_collection(name=collection_name)
    except Exception:
        # Fallback for older chromadb versions
        coll = ensure_collection(client, collection_name)

    embedder = Embedder(use_openai=True, openai_api_key=openai_api_key)

    total = 0
    for p in paths:
        n = process_single_file(p, course_id, lecture_id, embedder, coll, max_tokens=max_tokens, overlap=overlap)
        logger.info(f"Indexed {n} chunks from {p}")
        total += n

    # For persistent choma backends, call persist if supported
    try:
        client.persist()
    except Exception:
        pass

    logger.info(f"Total chunks indexed: {total}")
    return total


# ---------- CLI ----------
def find_input_paths(input_path: str) -> List[str]:
    p = Path(input_path)
    if p.is_dir():
        # collect supported files
        files = []
        for ext in ("*.pdf", "*.pptx", "*.md", "*.txt"):
            files.extend([str(x) for x in p.rglob(ext)])
        return sorted(files)
    elif p.is_file():
        return [str(p)]
    else:
        raise FileNotFoundError(f"Path not found: {input_path}")


def main():
    parser = argparse.ArgumentParser(description="Ingest course materials -> chunk -> embed -> Chroma")
    parser.add_argument("--input", "-i", required=True, help="File or directory to ingest")
    parser.add_argument("--collection", "-c", required=True, help="Chroma collection name (e.g., course_calculus)")
    parser.add_argument("--course-id", required=True, help="Course identifier")
    parser.add_argument("--lecture-id", required=False, help="Lecture identifier (optional)")
    parser.add_argument("--persist-dir", required=False, default="./chroma_db", help="Chroma persist directory")
    parser.add_argument("--openai-key", required=False, default=os.environ.get("OPENAI_API_KEY"), help="OpenAI API key (or set env OPENAI_API_KEY)")
    parser.add_argument("--max-tokens", type=int, default=500, help="Chunk size in tokens")
    parser.add_argument("--overlap", type=int, default=100, help="Chunk overlap in tokens")
    args = parser.parse_args()

    files = find_input_paths(args.input)
    if not files:
        print()
        logger.error("No files found to ingest.")
        return

    logger.info(f"Found {len(files)} files to ingest.")
    ingest_files(files, args.collection, args.course_id, args.lecture_id, args.persist_dir, args.openai_key, args.max_tokens, args.overlap)


if __name__ == "__main__":
    main()
